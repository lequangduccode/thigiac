# -*- coding: utf-8 -*-
"""Sinh slide thuyet trinh (PowerPoint) cho de tai nhan biet thit tuoi.

Tao bo slide 16:9 day du: tieu de, muc tieu, du lieu, phuong phap, ket qua
(kem % va hinh), demo web, han che, ket luan. Nhung hinh tu outputs/report_figures.

Usage: python scripts/build_slides.py
"""
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path(__file__).resolve().parents[1]
FIG = ROOT / "outputs" / "report_figures"
CM = ROOT / "outputs" / "locbeef_rf_v1" / "confusion_matrix.png"
OUT = ROOT / "reports" / "slide_thuyet_trinh.pptx"

# Bang mau
NAVY = RGBColor(0x0B, 0x25, 0x45)
BLUE = RGBColor(0x2E, 0x74, 0xB5)
GREEN = RGBColor(0x2F, 0xBF, 0x71)
RED = RGBColor(0xEF, 0x4D, 0x5B)
GREY = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xEE, 0xF3, 0xFA)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def _set(run, size, bold=False, color=None, italic=False, font="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    if color is not None:
        run.font.color.rgb = color


def textbox(slide, l, t, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if isinstance(ln, tuple):
            text, kw = ln
        else:
            text, kw = ln, {}
        p.space_after = Pt(kw.pop("space_after", 6))
        if kw.pop("bullet", False):
            text = "•  " + text
        run = p.add_run()
        run.text = text
        _set(run, kw.get("size", 20), kw.get("bold", False),
             kw.get("color", NAVY), kw.get("italic", False))
    return tb


def title_band(slide, title, idx=None):
    band = slide.shapes.add_shape(1, 0, 0, SW, Inches(1.15))
    band.fill.solid(); band.fill.fore_color.rgb = NAVY
    band.line.fill.background()
    tf = band.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.5)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    _set(r, 30, True, WHITE)
    if idx is not None:
        num = slide.shapes.add_textbox(SW - Inches(1.2), 0, Inches(1.0), Inches(1.15))
        ntf = num.text_frame; ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
        np_ = ntf.paragraphs[0]; np_.alignment = PP_ALIGN.RIGHT
        nr = np_.add_run(); nr.text = idx; _set(nr, 16, False, RGBColor(0x9A,0xB4,0xD0))


def fit_image(slide, path, box_l, box_t, box_w, box_h, caption=None):
    if not Path(path).exists():
        return
    iw, ih = Image.open(path).size
    ar = iw / ih
    box_ar = box_w / box_h
    if ar > box_ar:
        w = box_w; h = int(box_w / ar)
    else:
        h = box_h; w = int(box_h * ar)
    l = box_l + (box_w - w) // 2
    t = box_t + (box_h - h) // 2
    slide.shapes.add_picture(str(path), l, t, width=w, height=h)
    if caption:
        cap = slide.shapes.add_textbox(box_l, box_t + box_h + Inches(0.02), box_w, Inches(0.4))
        cp = cap.text_frame.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run(); cr.text = caption; _set(cr, 12, False, GREY, italic=True)


def table_slide(slide, l, t, headers, rows, col_w, row_h=Inches(0.5)):
    nrows, ncols = len(rows) + 1, len(headers)
    total_w = sum(col_w)
    tbl = slide.shapes.add_table(nrows, ncols, l, t, total_w, row_h * nrows).table
    for i, cw in enumerate(col_w):
        tbl.columns[i].width = cw
    for j, h in enumerate(headers):
        c = tbl.cell(0, j); c.fill.solid(); c.fill.fore_color.rgb = NAVY
        pr = c.text_frame.paragraphs[0]; pr.alignment = PP_ALIGN.CENTER
        rn = pr.add_run(); rn.text = h; _set(rn, 16, True, WHITE)
    for ri, row in enumerate(rows, 1):
        for j, val in enumerate(row):
            c = tbl.cell(ri, j)
            c.fill.solid(); c.fill.fore_color.rgb = LIGHT if ri % 2 else WHITE
            pr = c.text_frame.paragraphs[0]
            pr.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            rn = pr.add_run(); rn.text = val; _set(rn, 15, j == 0, NAVY)
    return tbl


# =====================================================================
# SLIDE 1 - TITLE
# =====================================================================
s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_shape(1, 0, 0, SW, SH)
bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
textbox(s, Inches(1), Inches(2.1), Inches(11.3), Inches(2.2), [
    ("NHẬN BIẾT THỊT TƯƠI", {"size": 46, "bold": True, "color": WHITE, "space_after": 6}),
    ("BẰNG THỊ GIÁC MÁY TÍNH", {"size": 46, "bold": True, "color": WHITE}),
], align=PP_ALIGN.CENTER)
textbox(s, Inches(1), Inches(4.3), Inches(11.3), Inches(1.6), [
    ("Phân loại độ tươi thịt bò từ ảnh bằng đặc trưng màu, texture và RandomForest",
     {"size": 22, "color": RGBColor(0xBF,0xD4,0xEC), "italic": True, "space_after": 14}),
    ("Bộ dữ liệu LocBeef (3.268 ảnh)  •  Accuracy 97,9%  •  Có ứng dụng web demo",
     {"size": 18, "color": GREEN, "bold": True}),
], align=PP_ALIGN.CENTER)
textbox(s, Inches(1), Inches(6.4), Inches(11.3), Inches(0.6),
        [("Học phần: Thị giác máy tính  |  Sinh viên thực hiện: (điền tên)",
          {"size": 15, "color": RGBColor(0x9A,0xB4,0xD0)})], align=PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 2 - NOI DUNG
# =====================================================================
s = prs.slides.add_slide(BLANK)
title_band(s, "Nội dung trình bày", "02")
textbox(s, Inches(1.2), Inches(1.6), Inches(11), Inches(5.5), [
    ("1.  Đặt vấn đề và mục tiêu", {"size": 24, "bold": True, "color": BLUE, "space_after": 12}),
    ("2.  Tổng quan phương pháp", {"size": 24, "bold": True, "color": BLUE, "space_after": 12}),
    ("3.  Bộ dữ liệu LocBeef", {"size": 24, "bold": True, "color": BLUE, "space_after": 12}),
    ("4.  Phương pháp đề xuất (pipeline, đặc trưng, mô hình)", {"size": 24, "bold": True, "color": BLUE, "space_after": 12}),
    ("5.  Kết quả thực nghiệm", {"size": 24, "bold": True, "color": BLUE, "space_after": 12}),
    ("6.  Demo ứng dụng web", {"size": 24, "bold": True, "color": BLUE, "space_after": 12}),
    ("7.  Hạn chế, bài học và kết luận", {"size": 24, "bold": True, "color": BLUE}),
])

# =====================================================================
# SLIDE 3 - DAT VAN DE
# =====================================================================
s = prs.slides.add_slide(BLANK)
title_band(s, "1. Đặt vấn đề", "03")
textbox(s, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.5), [
    ("Chất lượng thịt ảnh hưởng trực tiếp tới sức khỏe; thịt hỏng đổi màu và kết cấu bề mặt.", {"bullet": True, "space_after": 12}),
    ("Đánh giá bằng mắt mang tính chủ quan, thiếu nhất quán giữa người với người.", {"bullet": True, "space_after": 12}),
    ("Phương pháp dựa trên ảnh: không phá hủy mẫu, chi phí thấp, kết quả tức thì, dễ triển khai bằng camera.", {"bullet": True, "space_after": 12}),
    ("Mục tiêu: xây dựng hệ thống tự động phân loại ảnh thịt bò thành TƯƠI / HỎNG, kèm ứng dụng web.", {"bullet": True, "bold": True, "color": BLUE, "space_after": 12}),
    ("Ràng buộc: chạy trên CPU, dùng học máy cổ điển (không cần GPU/học sâu).", {"bullet": True, "space_after": 6}),
])

# =====================================================================
# SLIDE 4 - TONG QUAN
# =====================================================================
s = prs.slides.add_slide(BLANK)
title_band(s, "2. Tổng quan phương pháp", "04")
textbox(s, Inches(0.8), Inches(1.4), Inches(11.7), Inches(0.8), [
    ("Hai hướng tiếp cận chính cho bài toán đánh giá độ tươi thịt bằng ảnh:", {"size": 20, "space_after": 4})])
table_slide(s, Inches(1.2), Inches(2.3),
            ["Tiêu chí", "Đặc trưng thủ công + ML", "Học sâu (CNN)"],
            [
                ["Nhu cầu dữ liệu", "Vừa và nhỏ", "Lớn"],
                ["Tài nguyên", "CPU là đủ", "Cần GPU"],
                ["Khả năng giải thích", "Cao", "Thấp (hộp đen)"],
                ["Thời gian triển khai", "Nhanh", "Lâu hơn"],
            ],
            [Inches(3.6), Inches(4.2), Inches(3.2)])
textbox(s, Inches(1.2), Inches(6.2), Inches(11), Inches(0.9),
        [("→ Đề tài chọn hướng đặc trưng thủ công + RandomForest: nhẹ, dễ giải thích, phù hợp bài tập.",
          {"size": 19, "bold": True, "color": GREEN})])

# =====================================================================
# SLIDE 5 - DU LIEU
# =====================================================================
s = prs.slides.add_slide(BLANK)
title_band(s, "3. Bộ dữ liệu LocBeef", "05")
textbox(s, Inches(0.6), Inches(1.35), Inches(6.2), Inches(1.6), [
    ("3.268 ảnh thịt bò Aceh (Kaggle).", {"bullet": True, "size": 18, "space_after": 6}),
    ("2 lớp cân bằng: fresh 50% / spoiled 50%.", {"bullet": True, "size": 18, "space_after": 6}),
    ("Chia sẵn Train 70% / Test 30%.", {"bullet": True, "size": 18, "space_after": 6}),
])
table_slide(s, Inches(0.6), Inches(3.2),
            ["Tập", "fresh", "rotten", "Tổng", "Tỷ lệ"],
            [["Train", "1.144", "1.144", "2.288", "70%"],
             ["Test", "490", "490", "980", "30%"],
             ["Tổng", "1.634", "1.634", "3.268", "100%"]],
            [Inches(1.3), Inches(1.1), Inches(1.1), Inches(1.2), Inches(1.1)])
fit_image(s, FIG / "dataset_samples.png", Inches(7.0), Inches(1.5), Inches(6.0), Inches(5.2),
          "Hàng trên: tươi — Hàng dưới: hỏng")

# =====================================================================
# SLIDE 6 - PIPELINE
# =====================================================================
s = prs.slides.add_slide(BLANK)
title_band(s, "4. Phương pháp — Pipeline", "06")
fit_image(s, FIG / "pipeline_diagram.png", Inches(0.6), Inches(1.5), Inches(12.1), Inches(2.4))
textbox(s, Inches(0.8), Inches(4.2), Inches(11.7), Inches(2.8), [
    ("Resize 224×224 → cân bằng sáng CLAHE.", {"bullet": True, "space_after": 8}),
    ("Tách vùng thịt (loại nền trắng/sáng, vùng tối) để không học nhầm nền.", {"bullet": True, "space_after": 8}),
    ("Trích đặc trưng màu + texture 174 chiều trên vùng thịt.", {"bullet": True, "space_after": 8}),
    ("Phân loại bằng RandomForest (300 cây).", {"bullet": True}),
])

# =====================================================================
# SLIDE 7 - TACH VUNG THIT
# =====================================================================
s = prs.slides.add_slide(BLANK)
title_band(s, "4. Tách vùng thịt (masking)", "07")
fit_image(s, FIG / "preprocessing.png", Inches(0.6), Inches(1.5), Inches(12.1), Inches(3.6),
          "Ảnh gốc → CLAHE → mask vùng thịt → vùng thịt giữ lại")
textbox(s, Inches(0.8), Inches(5.7), Inches(11.7), Inches(1.4), [
    ("Loại nền trắng/đĩa và vùng quá tối; làm sạch bằng phép hình thái học (mở + đóng).", {"bullet": True, "space_after": 6}),
    ("Đặc trưng màu chỉ tính trên vùng thịt → giảm nhiễu nền, tăng độ chính xác.", {"bullet": True}),
])

# =====================================================================
# SLIDE 8 - DAC TRUNG
# =====================================================================
s = prs.slides.add_slide(BLANK)
title_band(s, "4. Đặc trưng 174 chiều", "08")
table_slide(s, Inches(1.4), Inches(1.7),
            ["Nhóm đặc trưng", "Chi tiết", "Số chiều", "Tỷ lệ"],
            [["Histogram HSV", "H32, S16, V16", "64", "36,8%"],
             ["Histogram Lab", "L16, a16, b16", "48", "27,6%"],
             ["Thống kê màu", "mean/std/p10/50/90 × 6", "30", "17,2%"],
             ["Texture LBP", "LBP 32 bins", "32", "18,4%"],
             ["Tổng cộng", "", "174", "100%"]],
            [Inches(3.0), Inches(3.8), Inches(1.7), Inches(1.7)])
textbox(s, Inches(1.4), Inches(6.3), Inches(11), Inches(0.9),
        [("Đặc trưng màu (HSV/Lab) mô tả độ đỏ/xỉn; LBP mô tả vân và độ thô mịn bề mặt.",
          {"size": 18, "italic": True, "color": GREY})])

# =====================================================================
# SLIDE 9 - KET QUA
# =====================================================================
s = prs.slides.add_slide(BLANK)
title_band(s, "5. Kết quả thực nghiệm", "09")
textbox(s, Inches(0.6), Inches(1.35), Inches(6.2), Inches(0.8),
        [("Accuracy trên test (980 ảnh): 97,9%", {"size": 24, "bold": True, "color": GREEN})])
table_slide(s, Inches(0.6), Inches(2.3),
            ["Chỉ số", "fresh", "spoiled"],
            [["Precision", "100,0%", "95,9%"],
             ["Recall", "95,7%", "100,0%"],
             ["F1-score", "97,8%", "97,9%"]],
            [Inches(2.2), Inches(1.8), Inches(1.8)])
textbox(s, Inches(0.6), Inches(5.0), Inches(6.2), Inches(1.9), [
    ("21 ảnh tươi bị nhầm thành hỏng; 0 ảnh hỏng bị nhầm thành tươi.", {"bullet": True, "size": 17, "space_after": 6}),
    ("Mô hình thiên về 'báo hỏng' → an toàn cho thực phẩm.", {"bullet": True, "size": 17, "color": BLUE, "bold": True}),
])
fit_image(s, CM, Inches(7.1), Inches(1.4), Inches(5.8), Inches(5.4))

# =====================================================================
# SLIDE 10 - FEATURE IMPORTANCE
# =====================================================================
s = prs.slides.add_slide(BLANK)
title_band(s, "5. Đặc trưng nào quyết định?", "10")
fit_image(s, FIG / "feature_importance.png", Inches(0.6), Inches(1.4), Inches(7.4), Inches(5.6))
textbox(s, Inches(8.2), Inches(1.9), Inches(4.7), Inches(4.5), [
    ("Nhóm màu đóng góp 96,4%:", {"size": 20, "bold": True, "color": NAVY, "space_after": 12}),
    ("HSV: 56,5%", {"bullet": True, "size": 19, "color": BLUE, "space_after": 8}),
    ("Thống kê màu: 27,4%", {"bullet": True, "size": 19, "color": BLUE, "space_after": 8}),
    ("Lab: 12,6%", {"bullet": True, "size": 19, "color": BLUE, "space_after": 8}),
    ("LBP texture: 3,6%", {"bullet": True, "size": 19, "color": GREY, "space_after": 14}),
    ("→ Màu sắc là dấu hiệu chính phân biệt tươi/hỏng.", {"size": 18, "bold": True, "color": GREEN}),
])

# =====================================================================
# SLIDE 11 - DEMO WEB
# =====================================================================
s = prs.slides.add_slide(BLANK)
title_band(s, "6. Demo ứng dụng web", "11")
fit_image(s, FIG / "web_result_fresh.png", Inches(0.5), Inches(1.5), Inches(6.1), Inches(4.6),
          "Ảnh thịt tươi → 'Tươi'")
fit_image(s, FIG / "web_result_spoiled.png", Inches(6.9), Inches(1.5), Inches(6.1), Inches(4.6),
          "Ảnh thịt hỏng → 'Hỏng'")
textbox(s, Inches(0.6), Inches(6.5), Inches(12), Inches(0.7),
        [("Tải ảnh lên trình duyệt → nhận nhãn + xác suất từng lớp; đọc được cả ảnh AVIF/HEIC.",
          {"size": 17, "italic": True, "color": GREY, "space_after": 0})], align=PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 12 - DEMO PREDICTIONS
# =====================================================================
s = prs.slides.add_slide(BLANK)
title_band(s, "6. Kết quả demo trên ảnh test", "12")
fit_image(s, FIG / "demo_predictions.png", Inches(1.4), Inches(1.4), Inches(10.5), Inches(5.5),
          "6/6 ảnh test được dự đoán đúng (✓)")

# =====================================================================
# SLIDE 13 - HAN CHE
# =====================================================================
s = prs.slides.add_slide(BLANK)
title_band(s, "7. Hạn chế và bài học", "13")
textbox(s, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.5), [
    ("Rò rỉ dữ liệu tiềm ẩn: dùng phân chia có sẵn — nếu cùng miếng thịt ở cả train/test, 97,9% có thể lạc quan.", {"bullet": True, "space_after": 14}),
    ("Lệch phân phối (domain shift): sai nhiều hơn với ảnh thịt ngẫu nhiên trên web (khác camera/ánh sáng).", {"bullet": True, "space_after": 14}),
    ("Bài học quan trọng: lớp 'hybrid' phân tích màu thủ công từng làm accuracy tụt còn 50% → đã loại bỏ.", {"bullet": True, "bold": True, "color": RED, "space_after": 14}),
    ("→ Heuristic cảm tính phải được kiểm chứng bằng đánh giá định lượng trên dữ liệu thật.", {"size": 19, "bold": True, "color": BLUE}),
])

# =====================================================================
# SLIDE 14 - KET LUAN
# =====================================================================
s = prs.slides.add_slide(BLANK)
title_band(s, "Kết luận và hướng phát triển", "14")
textbox(s, Inches(0.8), Inches(1.5), Inches(11.7), Inches(3.2), [
    ("Xây dựng pipeline: tách vùng thịt + đặc trưng 174 chiều + RandomForest.", {"bullet": True, "space_after": 10}),
    ("Đạt 97,9% accuracy trên test set thật LocBeef.", {"bullet": True, "space_after": 10}),
    ("Đóng gói thành ứng dụng web tải ảnh → xem kết quả.", {"bullet": True, "space_after": 10}),
])
textbox(s, Inches(0.8), Inches(4.6), Inches(11.7), Inches(2.4), [
    ("Hướng phát triển:", {"size": 22, "bold": True, "color": NAVY, "space_after": 10}),
    ("Chia dữ liệu theo mẫu vật; tăng đa dạng dữ liệu; transfer learning; cơ chế từ chối dự đoán khi thiếu tin cậy.",
     {"bullet": True, "size": 19, "color": BLUE}),
])

# =====================================================================
# SLIDE 15 - CAM ON
# =====================================================================
s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_shape(1, 0, 0, SW, SH)
bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
textbox(s, Inches(1), Inches(2.8), Inches(11.3), Inches(1.8), [
    ("CẢM ƠN THẦY CÔ VÀ CÁC BẠN ĐÃ LẮNG NGHE", {"size": 40, "bold": True, "color": WHITE, "space_after": 16}),
    ("Q & A", {"size": 28, "bold": True, "color": GREEN}),
], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print("Saved:", OUT, "| slides:", len(prs.slides._sldIdLst))
